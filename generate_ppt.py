#!/usr/bin/env python3
"""
W4Agent 本科毕设中期答辩PPT生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 颜色常量 ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BG = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xC7)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
LIGHT_BLUE = RGBColor(0x90, 0xE0, 0xEF)
LIGHT_BG = RGBColor(0xF0, 0xF7, 0xFF)
SECTION_BG = RGBColor(0xE8, 0xF4, 0xFD)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x55, 0x55, 0x55)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
TITLE_BLUE = RGBColor(0x0A, 0x2F, 0x5C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def add_bg(slide, color):
    """设置幻灯片纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    """添加一个矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加文字框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16, color=BLACK, line_spacing=1.5, bold=False, alignment=PP_ALIGN.LEFT):
    """添加多行文字"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=15, color=DARK_GRAY, bullet_color=ACCENT_BLUE, spacing=1.3):
    """添加带项目符号的列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # 使用圆点作为bullet
        run1 = p.add_run()
        run1.text = "●  "
        run1.font.size = Pt(font_size - 2)
        run1.font.color.rgb = bullet_color

        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = "Microsoft YaHei"
        p.space_after = Pt(font_size * (spacing - 1) + 2)
    return txBox


def add_slide_number(slide, num, total):
    """添加页码"""
    add_text(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.4),
             f"{num}/{total}", font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT)


def add_page_title(slide, title, subtitle=None):
    """添加页面标题栏 (蓝色横条)"""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), TITLE_BLUE)
    add_text(slide, Inches(0.7), Inches(0.15), Inches(10), Inches(0.7),
             title, font_size=30, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(0.7), Inches(10), Inches(0.4),
                 subtitle, font_size=14, color=LIGHT_BLUE)
    # 底部装饰线
    add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), ACCENT_TEAL)


def add_card(slide, left, top, width, height, title, content_lines, icon_text=None, title_color=TITLE_BLUE, bg=CARD_BG):
    """添加信息卡片"""
    card = add_rounded_rect(slide, left, top, width, height, bg, CARD_BORDER)
    # 顶部色条
    add_rect(slide, left + Inches(0.05), top + Inches(0.05), width - Inches(0.1), Inches(0.06), ACCENT_BLUE)

    y = top + Inches(0.2)
    if icon_text:
        add_text(slide, left + Inches(0.2), y, Inches(0.4), Inches(0.4),
                 icon_text, font_size=20, color=ACCENT_BLUE, bold=True)
        add_text(slide, left + Inches(0.6), y, width - Inches(0.8), Inches(0.35),
                 title, font_size=16, color=title_color, bold=True)
    else:
        add_text(slide, left + Inches(0.2), y, width - Inches(0.4), Inches(0.35),
                 title, font_size=16, color=title_color, bold=True)

    y += Inches(0.4)
    for line in content_lines:
        add_text(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.28),
                 line, font_size=12, color=MEDIUM_GRAY)
        y += Inches(0.26)
    return card


TOTAL_SLIDES = 14


# ============================================================
# SLIDE 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# 装饰元素
add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_BLUE)
add_rect(slide, Inches(0), Inches(3.3), SLIDE_W, Inches(0.04), ACCENT_TEAL)

# 标题
add_text(slide, Inches(1.0), Inches(1.2), Inches(11), Inches(0.6),
         "本科毕业设计中期答辩", font_size=16, color=LIGHT_BLUE, alignment=PP_ALIGN.LEFT)

add_multiline(slide, Inches(1.0), Inches(1.8), Inches(11), Inches(1.6), [
    "基于智能体技术的Web无障碍智能检测系统",
    "W4Agent — 从规则驱动到认知驱动"
], font_size=36, color=WHITE, bold=True, line_spacing=1.6)

# 分割线下方信息
info_lines = [
    "答辩人：XXX        学号：XXXXXXXXXX",
    "指导教师：XXX 教授",
    "专业：软件工程        院系：计算机科学与技术学院",
]
add_multiline(slide, Inches(1.0), Inches(4.0), Inches(8), Inches(1.8),
              info_lines, font_size=16, color=RGBColor(0xAA, 0xCC, 0xEE), line_spacing=2.0)

add_text(slide, Inches(1.0), Inches(6.5), Inches(5), Inches(0.4),
         "2025年 X月", font_size=14, color=GRAY)


# ============================================================
# SLIDE 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "目  录", "CONTENTS")

sections = [
    ("01", "研究背景与意义", "Web无障碍的现状与挑战"),
    ("02", "研究内容与目标", "关键科学问题与技术目标"),
    ("03", "系统总体设计", "架构设计与核心模块"),
    ("04", "多智能体协同引擎", "LangGraph状态机编排"),
    ("05", "自适应智能遍历器", "启发式探索与语义去重"),
    ("06", "多层次检测引擎", "规则+AI+视觉多模态融合"),
    ("07", "人机协同工作台", "实时监控与交互标注"),
    ("08", "中期进展与计划", "完成情况与后续安排"),
]

for i, (num, title, desc) in enumerate(sections):
    col = i % 2
    row = i // 2
    x = Inches(1.0) + col * Inches(5.8)
    y = Inches(1.5) + row * Inches(1.35)

    # 编号圆圈
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.05), Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE if i < 6 else ORANGE
    circle.line.fill.background()
    add_text(slide, x + Inches(0.05), y + Inches(0.1), Inches(0.45), Inches(0.45),
             num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, x + Inches(0.7), y + Inches(0.02), Inches(4.5), Inches(0.35),
             title, font_size=18, color=DARK_GRAY, bold=True)
    add_text(slide, x + Inches(0.7), y + Inches(0.38), Inches(4.5), Inches(0.3),
             desc, font_size=12, color=GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)


# ============================================================
# SLIDE 3: 研究背景与意义
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "01  研究背景与意义", "Research Background & Significance")

# 左侧: 背景
add_text(slide, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.4),
         "Web无障碍: 数字包容的基石", font_size=20, color=TITLE_BLUE, bold=True)

bg_items = [
    "全球超过13亿人（约16%）存在某种形式的残障 — WHO 2023",
    "WCAG 2.1/2.2 已成为国际Web无障碍核心标准",
    "我国《无障碍环境建设法》2023年正式施行",
    "96.3%的主流网站首页存在可检测的无障碍问题 — WebAIM 2024",
]
add_bullet_list(slide, Inches(0.7), Inches(1.9), Inches(5.5), Inches(2.5),
                bg_items, font_size=14, spacing=1.5)

# 右侧: 现有工具的局限
add_text(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.4),
         "现有自动化工具的局限", font_size=20, color=RED, bold=True)

problems = [
    ("遍历策略僵化", "固定规则爬取，无法适应动态内容和SPA"),
    ("干扰处理缺失", "弹窗/Cookie/广告阻断检测流程"),
    ("检测维度单一", "仅依赖DOM规则，难以发现语义级问题"),
    ("用户视角缺位", "无法从真实用户认知角度评估可用性"),
]
y = Inches(1.95)
for title, desc in problems:
    add_text(slide, Inches(7.2), y, Inches(0.2), Inches(0.25),
             "✕", font_size=14, color=RED, bold=True)
    add_text(slide, Inches(7.5), y, Inches(5), Inches(0.25),
             title, font_size=14, color=DARK_GRAY, bold=True)
    add_text(slide, Inches(7.5), y + Inches(0.28), Inches(5), Inches(0.25),
             desc, font_size=12, color=GRAY)
    y += Inches(0.6)

# 底部: 本研究意义
add_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2), SECTION_BG)
add_text(slide, Inches(0.8), Inches(4.9), Inches(3), Inches(0.4),
         "本研究的核心价值", font_size=18, color=TITLE_BLUE, bold=True)

values = [
    ("理论层面", "提出'认知驱动'的Web无障碍智能检测范式，融合多Agent协同与多模态感知"),
    ("技术层面", "构建端到端自动化系统，突破遍历、检测、响应三大技术瓶颈"),
    ("实践层面", "为残障用户提供更高质量的Web体验评估，助力数字包容社会建设"),
]
y = Inches(5.35)
for label, desc in values:
    add_text(slide, Inches(0.9), y, Inches(1.2), Inches(0.3),
             label, font_size=13, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(2.2), y, Inches(10.2), Inches(0.3),
             desc, font_size=13, color=DARK_GRAY)
    y += Inches(0.35)

add_slide_number(slide, 3, TOTAL_SLIDES)


# ============================================================
# SLIDE 4: 研究内容与目标
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "02  研究内容与目标", "Research Content & Objectives")

# 关键科学问题
add_text(slide, Inches(0.7), Inches(1.4), Inches(5), Inches(0.4),
         "关键科学问题", font_size=20, color=TITLE_BLUE, bold=True)

questions = [
    "Q1: 如何让多个智能体协同完成复杂的Web无障碍检测任务？",
    "Q2: 如何在动态Web环境中实现自适应、高覆盖率的页面遍历？",
    "Q3: 如何融合规则、语义与视觉多维度实现全面的无障碍问题检测？",
]
add_bullet_list(slide, Inches(0.7), Inches(1.9), Inches(5.8), Inches(1.5),
                questions, font_size=14, bullet_color=ORANGE, spacing=1.6)

# 技术目标
add_text(slide, Inches(7.0), Inches(1.4), Inches(5), Inches(0.4),
         "技术目标", font_size=20, color=TITLE_BLUE, bold=True)

goals = [
    ("多Agent编排引擎", "基于LangGraph实现Master/Crawler/\nDetector/Reporter四Agent协同"),
    ("自适应遍历器", "启发式评分+语义去重+突发事件响应"),
    ("多层检测引擎", "规则检测+AI语义+多模态视觉三层融合"),
    ("人机协同工作台", "实时监控+交互标注+报告生成"),
]
y = Inches(1.9)
for title, desc in goals:
    indicator = add_rounded_rect(slide, Inches(7.0), y, Inches(0.08), Inches(0.55), ACCENT_BLUE)
    add_text(slide, Inches(7.3), y, Inches(5.2), Inches(0.3),
             title, font_size=14, color=DARK_GRAY, bold=True)
    add_text(slide, Inches(7.3), y + Inches(0.3), Inches(5.2), Inches(0.3),
             desc, font_size=12, color=GRAY)
    y += Inches(0.65)

# 创新点
add_rect(slide, Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5), SECTION_BG)
add_text(slide, Inches(0.8), Inches(4.7), Inches(3), Inches(0.4),
         "创新点", font_size=18, color=TITLE_BLUE, bold=True)

innovations = [
    ("创新点一", "基于LangGraph的多智能体动态编排", "Master Agent根据检测进度动态调度探索/检测/报告阶段，支持自适应迭代决策"),
    ("创新点二", "启发式+语义的自适应遍历策略", "结合URL特征评分和DOM结构模板签名去重，解决动态页面高效覆盖问题"),
    ("创新点三", "规则-语义-视觉三层融合检测", "首创将多模态LLM视觉分析引入无障碍检测，对比元素外观与标签语义一致性"),
    ("创新点四", "长短期记忆驱动的跨会话学习", "结合情景/语义/程序性三类记忆，使系统能在同域检测中持续积累经验"),
]
y = Inches(5.15)
for num, title, desc in innovations:
    add_text(slide, Inches(0.9), y, Inches(0.9), Inches(0.28),
             num, font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(1.8), y, Inches(2.2), Inches(0.28),
             title, font_size=12, color=DARK_GRAY, bold=True)
    add_text(slide, Inches(4.0), y, Inches(8.5), Inches(0.28),
             desc, font_size=11, color=MEDIUM_GRAY)
    y += Inches(0.35)

add_slide_number(slide, 4, TOTAL_SLIDES)


# ============================================================
# SLIDE 5: 系统总体架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "03  系统总体设计", "System Architecture Overview")

# 技术栈概览
add_text(slide, Inches(0.7), Inches(1.35), Inches(4), Inches(0.35),
         "技术栈", font_size=18, color=TITLE_BLUE, bold=True)

tech_items = [
    ("后端", "Python 3.11 + FastAPI + LangChain + LangGraph + Playwright"),
    ("前端", "React 18 + TypeScript + Ant Design + Zustand + Vite"),
    ("数据库", "PostgreSQL (pgvector) + Redis"),
    ("LLM", "OpenAI / Anthropic Claude / 本地Ollama (可配置)"),
    ("部署", "Docker Compose 四服务编排"),
]
y = Inches(1.75)
for label, value in tech_items:
    add_text(slide, Inches(0.8), y, Inches(0.8), Inches(0.25),
             label, font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(1.65), y, Inches(4.5), Inches(0.25),
             value, font_size=12, color=DARK_GRAY)
    y += Inches(0.3)

# 架构图（文字版）
add_text(slide, Inches(0.7), Inches(3.3), Inches(4), Inches(0.35),
         "系统架构 (四层)", font_size=18, color=TITLE_BLUE, bold=True)

layers = [
    ("表现层", "React SPA + WebSocket实时通信", ACCENT_BLUE),
    ("服务层", "FastAPI REST API + 业务逻辑层", GREEN),
    ("引擎层", "Multi-Agent Engine + Adaptive Crawler + Detection Engine", ORANGE),
    ("数据层", "PostgreSQL + Redis + 文件存储", PURPLE),
]
y = Inches(3.8)
for name, desc, color in layers:
    add_rounded_rect(slide, Inches(0.8), y, Inches(5.2), Inches(0.55), color)
    add_text(slide, Inches(1.0), y + Inches(0.02), Inches(1.2), Inches(0.28),
             name, font_size=13, color=WHITE, bold=True)
    add_text(slide, Inches(1.0), y + Inches(0.28), Inches(4.8), Inches(0.25),
             desc, font_size=11, color=WHITE)
    y += Inches(0.65)

# 右侧：四大核心模块
add_text(slide, Inches(6.8), Inches(1.35), Inches(5), Inches(0.35),
         "四大核心模块", font_size=18, color=TITLE_BLUE, bold=True)

modules = [
    ("多智能体编排引擎", [
        "Master/Crawler/Detector/Reporter",
        "LangGraph状态图动态编排",
        "长短期记忆系统",
    ], ACCENT_BLUE),
    ("自适应遍历器", [
        "启发式URL优先级评分",
        "DOM结构模板语义去重",
        "突发事件自动响应",
    ], GREEN),
    ("多层次检测引擎", [
        "7条WCAG规则自动检查",
        "AI语义分析补充检测",
        "多模态视觉对比分析",
    ], ORANGE),
    ("人机协同工作台", [
        "实时WebSocket监控",
        "可视化标注截图",
        "报告生成与导出",
    ], PURPLE),
]
y = Inches(1.8)
for title, items, color in modules:
    add_rounded_rect(slide, Inches(6.8), y, Inches(5.6), Inches(1.1), LIGHT_BG, color)
    add_text(slide, Inches(7.0), y + Inches(0.05), Inches(5), Inches(0.3),
             title, font_size=14, color=color, bold=True)
    item_y = y + Inches(0.35)
    for item in items:
        add_text(slide, Inches(7.2), item_y, Inches(5), Inches(0.22),
                 f"• {item}", font_size=11, color=MEDIUM_GRAY)
        item_y += Inches(0.22)
    y += Inches(1.2)

add_slide_number(slide, 5, TOTAL_SLIDES)


# ============================================================
# SLIDE 6: 数据库设计
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "03  数据库与接口设计", "Database & API Design")

# 数据库模型
add_text(slide, Inches(0.7), Inches(1.4), Inches(5), Inches(0.35),
         "数据库模型 (7张表)", font_size=18, color=TITLE_BLUE, bold=True)

tables = [
    ("projects", "项目管理", "name, description, base_url, owner"),
    ("tasks", "检测任务", "target_url, status(6态), config(JSON), progress"),
    ("pages", "页面数据", "url, title, depth, content_hash, a11y_tree, screenshot"),
    ("issues", "问题记录", "wcag_criterion, severity, detected_by, confidence, bbox"),
    ("reports", "检测报告", "scores, severity_counts, summary, recommendations"),
    ("annotations", "人工标注", "type(confirm/false_positive/reclassify/comment)"),
    ("agent_memories", "Agent记忆", "type(episodic/semantic/procedural), domain, content"),
]
y = Inches(1.85)
for name, desc, fields in tables:
    add_text(slide, Inches(0.8), y, Inches(1.6), Inches(0.25),
             name, font_size=12, color=ACCENT_BLUE, bold=True)
    add_text(slide, Inches(2.5), y, Inches(1.2), Inches(0.25),
             desc, font_size=12, color=DARK_GRAY, bold=True)
    add_text(slide, Inches(3.8), y, Inches(4), Inches(0.25),
             fields, font_size=11, color=GRAY)
    y += Inches(0.35)

# 右侧: API接口设计
add_text(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.35),
         "REST API (6组路由) + WebSocket", font_size=18, color=TITLE_BLUE, bold=True)

apis = [
    ("项目管理", "CRUD操作, 分页查询"),
    ("任务管理", "创建/启动/停止/查询任务"),
    ("页面数据", "按任务获取页面列表及详情"),
    ("问题管理", "按任务/页面/严重度筛选, 标注"),
    ("报告接口", "获取报告, 导出HTML/PDF/JSON"),
    ("系统设置", "LLM配置, WCAG级别配置"),
    ("WebSocket", "实时推送: 进度/发现/问题/推理日志"),
]
y = Inches(1.85)
for name, desc in apis:
    add_text(slide, Inches(7.3), y, Inches(1.6), Inches(0.25),
             name, font_size=12, color=GREEN, bold=True)
    add_text(slide, Inches(9.0), y, Inches(3.8), Inches(0.25),
             desc, font_size=12, color=DARK_GRAY)
    y += Inches(0.35)

# 底部特色说明
add_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), SECTION_BG)
add_text(slide, Inches(0.8), Inches(5.1), Inches(10), Inches(0.35),
         "设计特色", font_size=16, color=TITLE_BLUE, bold=True)

features_text = [
    "UUID主键 + 自动时间戳：确保分布式环境下的唯一性与可追溯性",
    "JSON字段灵活存储：task.config、page.a11y_tree、issue.context 支持动态扩展",
    "六状态任务生命周期：PENDING → RUNNING → PAUSED → COMPLETED / FAILED / CANCELLED",
    "三来源问题检测：detected_by 区分 rule / ai / vision_ai，支持混合检测结果分析",
    "WebSocket实时推送：每个检测事件即时通知前端，无需轮询",
]
add_bullet_list(slide, Inches(0.7), Inches(5.45), Inches(11.5), Inches(1.5),
                features_text, font_size=12, spacing=1.4)

add_slide_number(slide, 6, TOTAL_SLIDES)


# ============================================================
# SLIDE 7: 多智能体协同引擎
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "04  多智能体协同引擎", "Multi-Agent Orchestration Engine")

# 状态机流程
add_text(slide, Inches(0.7), Inches(1.4), Inches(6), Inches(0.35),
         "LangGraph 状态机编排", font_size=20, color=TITLE_BLUE, bold=True)

# 绘制流程图（简化版，用形状）
flow_y = Inches(2.0)
flow_nodes = [
    (Inches(1.0), "START", GRAY),
    (Inches(3.0), "Plan\n(Master)", ACCENT_BLUE),
    (Inches(5.5), "Explore\n(Crawler)", GREEN),
    (Inches(8.0), "Test\n(Detector)", ORANGE),
    (Inches(10.5), "Report\n(Reporter)", PURPLE),
]

for x, text, color in flow_nodes:
    if text == "START":
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, flow_y + Inches(0.1), Inches(1.0), Inches(0.6))
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, flow_y, Inches(1.5), Inches(0.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# 箭头用文字表示
arrows = [
    (Inches(2.05), flow_y + Inches(0.3), "→"),
    (Inches(4.55), flow_y + Inches(0.3), "→"),
    (Inches(7.05), flow_y + Inches(0.3), "→"),
    (Inches(9.55), flow_y + Inches(0.3), "→"),
]
for x, y_pos, text in arrows:
    add_text(slide, x, y_pos, Inches(0.5), Inches(0.3),
             text, font_size=24, color=DARK_GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# 循环箭头标注
add_text(slide, Inches(4.2), flow_y + Inches(0.85), Inches(5), Inches(0.3),
         "← ← ←  Master 根据进度动态决策循环  → → →", font_size=11, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Pipeline State
add_text(slide, Inches(0.7), Inches(3.3), Inches(5.5), Inches(0.35),
         "Pipeline 共享状态 (PipelineState)", font_size=16, color=TITLE_BLUE, bold=True)

state_items = [
    "target_url: 目标网站入口地址",
    "explored_urls / pending_urls: 已探索/待探索URL集合",
    "pending_test_urls: 待检测页面队列",
    "all_issues: 累计发现的问题列表",
    "iteration / max_iterations: 迭代计数与安全上限",
    "report_data: 最终报告数据",
]
add_bullet_list(slide, Inches(0.7), Inches(3.7), Inches(5.5), Inches(2.0),
                state_items, font_size=12, spacing=1.3)

# 右侧: 四个Agent详情
add_text(slide, Inches(6.8), Inches(3.3), Inches(5.5), Inches(0.35),
         "四大Agent职责", font_size=16, color=TITLE_BLUE, bold=True)

agents_info = [
    ("Master Agent", "全局调度决策", "基于进度统计决定下一步:EXPLORE/TEST/REPORT/COMPLETE\n结合领域长期记忆优化策略", ACCENT_BLUE),
    ("Crawler Agent", "智能导航决策", "分析A11y树与交互元素，决定点击/导航/滚动/填写\n智能截断输入(50元素,20URL)控制Token", GREEN),
    ("Detector Agent", "多维度检测", "三层分析:规则检查→AI语义→多模态视觉\n视觉分析对比元素外观与标签语义一致性", ORANGE),
    ("Reporter Agent", "报告生成", "汇总统计数据生成摘要与改进建议\n自动计算各WCAG级别达标分数", PURPLE),
]
y = Inches(3.7)
for name, role, desc, color in agents_info:
    add_rounded_rect(slide, Inches(6.8), y, Inches(5.8), Inches(0.85), LIGHT_BG, color)
    add_text(slide, Inches(7.0), y + Inches(0.02), Inches(2), Inches(0.28),
             name, font_size=13, color=color, bold=True)
    add_text(slide, Inches(9.0), y + Inches(0.02), Inches(3.4), Inches(0.28),
             role, font_size=12, color=DARK_GRAY, bold=True)
    # 描述拆分为两行
    desc_lines = desc.split('\n')
    for j, dl in enumerate(desc_lines):
        add_text(slide, Inches(7.0), y + Inches(0.3 + j * 0.22), Inches(5.4), Inches(0.22),
                 dl, font_size=10, color=MEDIUM_GRAY)
    y += Inches(0.92)

add_slide_number(slide, 7, TOTAL_SLIDES)


# ============================================================
# SLIDE 8: 记忆系统
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "04  Agent记忆系统", "Short-term & Long-term Memory")

# 短期记忆
add_text(slide, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.35),
         "短期记忆 (Short-term Memory)", font_size=20, color=ACCENT_BLUE, bold=True)

stm_items = [
    "基于有界双端队列 (deque)，容量100条",
    "记录类型: 发现页面/检测问题/Agent推理/遍历动作",
    "提供 get_context_summary() 生成LLM上下文",
    "单次检测任务内共享，任务结束后释放",
    "作用: 为Agent提供当前任务的即时上下文",
]
add_bullet_list(slide, Inches(0.7), Inches(1.9), Inches(5.5), Inches(2.0),
                stm_items, font_size=14, spacing=1.4)

# 长期记忆
add_text(slide, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.35),
         "长期记忆 (Long-term Memory)", font_size=20, color=PURPLE, bold=True)

ltm_items = [
    "持久化存储于PostgreSQL (agent_memories表)",
    "三类记忆类型，模仿人类认知结构：",
]
add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(0.8),
                ltm_items, font_size=14, spacing=1.3, bullet_color=PURPLE)

memory_types = [
    ("情景记忆\n(Episodic)", "记录具体检测经历\n\"该域名XX页面发现\nN个对比度问题\"", ACCENT_BLUE),
    ("语义记忆\n(Semantic)", "总结领域模式知识\n\"电商网站通常在商品\n详情页缺少图片替代文本\"", GREEN),
    ("程序性记忆\n(Procedural)", "积累处理策略经验\n\"SPA页面需等待路由\n加载完成后再检测\"", ORANGE),
]
x_start = Inches(7.0)
for i, (title, desc, color) in enumerate(memory_types):
    x = x_start + i * Inches(1.95)
    add_rounded_rect(slide, x, Inches(3.2), Inches(1.8), Inches(2.2), LIGHT_BG, color)
    add_text(slide, x + Inches(0.1), Inches(3.3), Inches(1.6), Inches(0.55),
             title, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(3.9), Inches(1.6), Inches(1.3),
             desc, font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# 底部: 记忆流程
add_rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.5), SECTION_BG)
add_text(slide, Inches(0.8), Inches(5.7), Inches(10), Inches(0.35),
         "记忆工作流程", font_size=16, color=TITLE_BLUE, bold=True)

flow_items = [
    "1. 任务启动时 → recall_for_context() 加载当前域名相关的长期记忆作为LLM上下文",
    "2. 检测过程中 → 重要发现实时写入短期记忆，为后续Agent决策提供上下文",
    "3. 任务完成后 → 将有价值的检测经验从短期记忆提炼并持久化为长期记忆",
    "4. 跨任务复用 → 同域名再次检测时，自动加载历史经验，实现持续学习与优化",
]
add_bullet_list(slide, Inches(0.7), Inches(6.05), Inches(11.5), Inches(1.0),
                flow_items, font_size=12, spacing=1.3, bullet_color=PURPLE)

add_slide_number(slide, 8, TOTAL_SLIDES)


# ============================================================
# SLIDE 9: 自适应遍历器
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "05  自适应智能遍历器", "Adaptive Intelligent Crawler")

# 左侧: 遍历流程
add_text(slide, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.35),
         "核心遍历流程", font_size=20, color=TITLE_BLUE, bold=True)

crawl_steps = [
    "1. 导航到目标URL → 等待页面完全加载",
    "2. 突发事件响应 → 自动处理弹窗/Cookie/广告",
    "3. 页面深度分析 → 提取A11y树/DOM/交互元素/Bbox",
    "4. 语义去重检查 → 内容哈希 + 结构模板匹配",
    "5. 状态图更新 → 记录节点与边的状态流转",
    "6. 全页面截图 → 保存用于后续视觉分析",
    "7. 链接发现 → 同域链接提取与去重",
    "8. 启发式排序 → 计算优先级评分，入队待探索",
]
add_bullet_list(slide, Inches(0.7), Inches(1.85), Inches(5.5), Inches(3.0),
                crawl_steps, font_size=13, spacing=1.35)

# 右侧: 三大子模块
add_text(slide, Inches(6.8), Inches(1.4), Inches(5.5), Inches(0.35),
         "三大核心机制", font_size=20, color=TITLE_BLUE, bold=True)

# 启发式评分
add_card(slide, Inches(6.8), Inches(1.85), Inches(5.6), Inches(1.55),
         "启发式URL优先级评分", [
             "基础分50分, 总分范围0-100",
             "高优先: 登录/注册/表单/搜索 (+30)",
             "深度惩罚: 随depth/max_depth线性递减",
             "父页面加分: 含表单+15, 有问题+10",
         ], icon_text="01")

# 语义去重
add_card(slide, Inches(6.8), Inches(3.55), Inches(5.6), Inches(1.55),
         "两层语义去重", [
             "精确层: SHA-256内容哈希直接比对",
             "模糊层: DOM结构签名(标题/表单/地标/图片)",
             "模板检测: 同模板3+实例后自动跳过",
             "有效避免商品页/博客等模板化页面重复",
         ], icon_text="02")

# 事件响应
add_card(slide, Inches(6.8), Inches(5.25), Inches(5.6), Inches(1.55),
         "突发事件响应器", [
             "检测4类干扰: Cookie同意/弹窗/广告/Alert",
             "三级处理: 专用按钮 → 通用关闭 → Escape键",
             "CSS选择器模式匹配 + 可见性验证",
             "所有处理事件实时WebSocket通知",
         ], icon_text="03")

add_slide_number(slide, 9, TOTAL_SLIDES)


# ============================================================
# SLIDE 10: 多层次检测引擎
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "06  多层次检测引擎", "Multi-layer Detection Engine")

# 三层检测架构
add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.35),
         "三层融合检测架构: 规则层 → AI语义层 → 多模态视觉层", font_size=20, color=TITLE_BLUE, bold=True)

# Layer 1: 规则检测
add_rounded_rect(slide, Inches(0.5), Inches(1.95), Inches(3.9), Inches(3.4), LIGHT_BG, ACCENT_BLUE)
add_text(slide, Inches(0.7), Inches(2.05), Inches(3.5), Inches(0.35),
         "Layer 1: 规则检测", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.7), Inches(2.4), Inches(3.5), Inches(0.3),
         "确定性WCAG规则检查", font_size=12, color=GRAY)

rules = [
    "img-alt: 图片替代文本缺失",
    "img-alt-quality: 无意义替代文本",
    "form-label: 表单标签关联",
    "heading-order: 标题层级跳跃",
    "html-lang: 页面语言声明",
    "landmark: 地标结构完整性",
    "link-text: 链接文本可理解性",
]
y = Inches(2.75)
for rule in rules:
    add_text(slide, Inches(0.8), y, Inches(3.4), Inches(0.22),
             f"• {rule}", font_size=10, color=MEDIUM_GRAY)
    y += Inches(0.24)

add_text(slide, Inches(0.7), y + Inches(0.05), Inches(3.5), Inches(0.25),
         "可扩展: RuleRegistry注册机制", font_size=10, color=ACCENT_BLUE)

# Layer 2: AI语义检测
add_rounded_rect(slide, Inches(4.7), Inches(1.95), Inches(3.9), Inches(3.4), LIGHT_BG, GREEN)
add_text(slide, Inches(4.9), Inches(2.05), Inches(3.5), Inches(0.35),
         "Layer 2: AI语义分析", font_size=16, color=GREEN, bold=True)
add_text(slide, Inches(4.9), Inches(2.4), Inches(3.5), Inches(0.3),
         "LLM深度语义理解", font_size=12, color=GRAY)

ai_items = [
    "输入A11y树+HTML片段+规则结果",
    "覆盖WCAG四大原则:",
    "  可感知: 替代文本质量/颜色对比",
    "  可操作: 键盘可访问/焦点管理",
    "  可理解: 标签清晰/错误提示",
    "  鲁棒性: ARIA正确使用",
    "输出结构化JSON问题列表",
    "与规则结果互补，不重复计算",
]
y = Inches(2.75)
for item in ai_items:
    add_text(slide, Inches(5.0), y, Inches(3.4), Inches(0.22),
             f"• {item}" if not item.startswith("  ") else f"  {item}", font_size=10, color=MEDIUM_GRAY)
    y += Inches(0.24)

# Layer 3: 视觉分析
add_rounded_rect(slide, Inches(8.9), Inches(1.95), Inches(3.9), Inches(3.4), LIGHT_BG, ORANGE)
add_text(slide, Inches(9.1), Inches(2.05), Inches(3.5), Inches(0.35),
         "Layer 3: 多模态视觉分析", font_size=16, color=ORANGE, bold=True)
add_text(slide, Inches(9.1), Inches(2.4), Inches(3.5), Inches(0.3),
         "视觉LLM看图识别问题", font_size=12, color=GRAY)

vision_items = [
    "截图编码为Base64输入视觉LLM",
    "附加元素Bounding Box坐标",
    "核心检测维度:",
    "  \"元素看起来像什么\"",
    "  vs \"标签说的是什么\"",
    "  → 捕获语义不一致问题",
    "自动与已有问题去重",
    "发现纯规则/文本无法检测的问题",
]
y = Inches(2.75)
for item in vision_items:
    add_text(slide, Inches(9.2), y, Inches(3.4), Inches(0.22),
             f"• {item}" if not item.startswith("  ") else f"  {item}", font_size=10, color=MEDIUM_GRAY)
    y += Inches(0.24)

# 底部: 可视化标注
add_rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.5), SECTION_BG)
add_text(slide, Inches(0.8), Inches(5.7), Inches(10), Inches(0.35),
         "检测结果可视化: VisualAnnotator", font_size=16, color=TITLE_BLUE, bold=True)

annotator_items = [
    "四级严重度色彩编码: 严重(红) / 重要(橙) / 轻微(金) / 提示(蓝)",
    "四层匹配策略: 精确CSS选择器 → HTML属性匹配 → 模糊标签+属性 → 位置回退",
    "自动生成带标注的截图 + 图例说明，直观展示问题位置和类型",
    "支持CJK中文字体渲染，适配Linux/macOS双平台",
]
add_bullet_list(slide, Inches(0.7), Inches(6.1), Inches(11.5), Inches(0.9),
                annotator_items, font_size=12, spacing=1.3, bullet_color=ORANGE)

add_slide_number(slide, 10, TOTAL_SLIDES)


# ============================================================
# SLIDE 11: 人机协同工作台
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "07  人机协同工作台", "Human-AI Collaborative Workspace")

# 前端技术栈
add_text(slide, Inches(0.7), Inches(1.4), Inches(5.5), Inches(0.35),
         "前端技术方案", font_size=18, color=TITLE_BLUE, bold=True)

fe_items = [
    "React 18 + TypeScript 5.5: 类型安全的现代前端框架",
    "Ant Design 5.20: 企业级UI组件库",
    "Zustand 4.5: 轻量级状态管理",
    "Vite 5.4: 极速开发构建工具",
    "WebSocket: 实时双向通信",
]
add_bullet_list(slide, Inches(0.7), Inches(1.8), Inches(5.5), Inches(1.8),
                fe_items, font_size=13, spacing=1.3)

# 六大页面
add_text(slide, Inches(6.8), Inches(1.4), Inches(5.5), Inches(0.35),
         "核心功能页面", font_size=18, color=TITLE_BLUE, bold=True)

pages_info = [
    ("Dashboard", "全局总览: 项目统计/最近任务/问题趋势"),
    ("项目管理", "项目CRUD/URL配置/检测历史"),
    ("任务创建", "配置目标URL/WCAG级别/最大页面数"),
    ("任务详情", "实时监控: 进度/问题/推理日志/截图(4 Tab)"),
    ("报告查看", "评分环形图/问题分布/标注截图/导出"),
    ("系统设置", "LLM选择/API Key/检测参数配置"),
]
y = Inches(1.8)
for name, desc in pages_info:
    add_text(slide, Inches(7.0), y, Inches(1.5), Inches(0.25),
             name, font_size=12, color=GREEN, bold=True)
    add_text(slide, Inches(8.6), y, Inches(4), Inches(0.25),
             desc, font_size=12, color=DARK_GRAY)
    y += Inches(0.32)

# 底部: 关键交互特性
add_rect(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.9), SECTION_BG)
add_text(slide, Inches(0.8), Inches(4.3), Inches(10), Inches(0.35),
         "关键交互特性", font_size=18, color=TITLE_BLUE, bold=True)

# 三栏布局
features = [
    ("实时监控", [
        "WebSocket推送检测事件",
        "进度条实时更新",
        "Agent推理过程时间线",
        "指数退避自动重连(5次)",
    ], ACCENT_BLUE),
    ("可视化标注", [
        "原始/标注截图一键切换",
        "悬浮显示问题详情",
        "Bounding Box交互覆盖",
        "严重度颜色编码",
    ], GREEN),
    ("报告与导出", [
        "WCAG达标评分环形图",
        "问题按严重度/类型统计",
        "AI生成的摘要与建议",
        "支持HTML/PDF/JSON导出",
    ], ORANGE),
]

for i, (title, items, color) in enumerate(features):
    x = Inches(0.7) + i * Inches(4.1)
    add_rounded_rect(slide, x, Inches(4.75), Inches(3.8), Inches(2.2), CARD_BG, color)
    add_text(slide, x + Inches(0.15), Inches(4.85), Inches(3.5), Inches(0.3),
             title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    item_y = Inches(5.2)
    for item in items:
        add_text(slide, x + Inches(0.2), item_y, Inches(3.4), Inches(0.22),
                 f"• {item}", font_size=11, color=MEDIUM_GRAY)
        item_y += Inches(0.25)

add_slide_number(slide, 11, TOTAL_SLIDES)


# ============================================================
# SLIDE 12: 中期进展
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "08  中期进展", "Mid-term Progress")

# 已完成工作
add_text(slide, Inches(0.7), Inches(1.4), Inches(12), Inches(0.35),
         "已完成工作", font_size=20, color=GREEN, bold=True)

done_items = [
    ("系统架构设计", "完成四层架构设计、数据库ER图、API接口规范", "100%"),
    ("多Agent编排引擎", "LangGraph状态图+4个Agent+记忆系统完整实现", "100%"),
    ("自适应遍历器", "启发式评分+语义去重+事件响应+状态图全部实现", "100%"),
    ("规则检测引擎", "7条WCAG规则+Registry可扩展架构", "100%"),
    ("AI语义检测", "Detector Agent文本分析+Prompt设计", "100%"),
    ("多模态视觉分析", "截图Base64+Bbox+视觉LLM+结果去重", "100%"),
    ("可视化标注器", "四层Bbox匹配+严重度色彩+图例生成", "100%"),
    ("前端工作台", "6个页面+实时WebSocket+标注截图组件", "100%"),
    ("Docker部署", "四服务Compose编排+Nginx反向代理", "100%"),
]

y = Inches(1.85)
for item, desc, pct in done_items:
    # 进度条背景
    add_rounded_rect(slide, Inches(10.8), y + Inches(0.02), Inches(1.5), Inches(0.22), RGBColor(0xE0, 0xE0, 0xE0))
    # 进度条填充
    bar_width = 1.5 * int(pct.rstrip('%')) / 100
    add_rounded_rect(slide, Inches(10.8), y + Inches(0.02), Inches(bar_width), Inches(0.22), GREEN)
    add_text(slide, Inches(11.0), y + Inches(0.01), Inches(1.0), Inches(0.22),
             pct, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.8), y, Inches(2.5), Inches(0.25),
             f"✓  {item}", font_size=12, color=GREEN, bold=True)
    add_text(slide, Inches(3.5), y, Inches(7.2), Inches(0.25),
             desc, font_size=11, color=DARK_GRAY)
    y += Inches(0.32)

# 代码统计
add_rect(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.7), SECTION_BG)
add_text(slide, Inches(0.8), Inches(4.95), Inches(12), Inches(0.3),
         "代码规模:  后端 ~4,500行 Python  |  前端 ~3,000行 TypeScript/TSX  |  共计 ~7,500行  |  7张数据表  |  6组API + WebSocket",
         font_size=13, color=TITLE_BLUE, bold=True, alignment=PP_ALIGN.LEFT)

# 遇到的困难与解决
add_text(slide, Inches(0.7), Inches(5.8), Inches(12), Inches(0.35),
         "主要技术难点与解决方案", font_size=18, color=ORANGE, bold=True)

challenges = [
    ("HTTP/2代理兼容", "LLM API调用挂起 → 自定义httpx客户端强制HTTP/1.1"),
    ("截图Bbox对齐", "滚动位置偏移 → 使用page-absolute坐标+全页面截图"),
    ("Token控制", "Agent输入过长 → 智能截断(50元素, 20URL, 50KB HTML)"),
    ("模板页去重", "商品页重复遍历 → DOM结构签名+3实例阈值检测"),
]
y = Inches(6.2)
for challenge, solution in challenges:
    add_text(slide, Inches(0.9), y, Inches(2.5), Inches(0.22),
             challenge, font_size=11, color=RED, bold=True)
    add_text(slide, Inches(3.5), y, Inches(9), Inches(0.22),
             solution, font_size=11, color=DARK_GRAY)
    y += Inches(0.28)

add_slide_number(slide, 12, TOTAL_SLIDES)


# ============================================================
# SLIDE 13: 后续计划
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, "08  后续工作计划", "Future Work Plan")

# 甘特图式时间线
add_text(slide, Inches(0.7), Inches(1.4), Inches(10), Inches(0.35),
         "剩余工作时间规划", font_size=20, color=TITLE_BLUE, bold=True)

# 时间轴
months = ["4月", "5月上", "5月下", "6月"]
for i, m in enumerate(months):
    x = Inches(4.5) + i * Inches(2.1)
    add_rect(slide, x, Inches(1.9), Inches(2.0), Inches(0.35), ACCENT_BLUE)
    add_text(slide, x, Inches(1.92), Inches(2.0), Inches(0.3),
             m, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

tasks_plan = [
    ("系统测试与评估", Inches(4.5), Inches(4.2), [
        "对标axe-core/WAVE等工具进行对比实验",
        "检测覆盖率/准确率/误报率量化评估",
        "多网站类型(政府/电商/教育)测试",
    ]),
    ("性能优化", Inches(6.6), Inches(2.1), [
        "并发检测性能优化",
        "LLM调用成本控制",
    ]),
    ("论文撰写", Inches(6.6), Inches(4.2), [
        "研究方法与实验结果章节",
        "论文修改与定稿",
    ]),
    ("答辩准备", Inches(10.7), Inches(2.1), [
        "PPT制作与演示准备",
        "系统Demo演示",
    ]),
]

y = Inches(2.5)
colors = [GREEN, ORANGE, PURPLE, RED]
for i, (name, start_x, width, items) in enumerate(tasks_plan):
    # 任务标签
    add_text(slide, Inches(0.8), y, Inches(3.5), Inches(0.3),
             name, font_size=13, color=colors[i], bold=True)
    # 甘特条
    add_rounded_rect(slide, start_x, y + Inches(0.03), width, Inches(0.28), colors[i])
    # 详情
    item_y = y + Inches(0.35)
    for item in items:
        add_text(slide, Inches(0.9), item_y, Inches(3.3), Inches(0.22),
                 f"• {item}", font_size=10, color=MEDIUM_GRAY)
        item_y += Inches(0.22)
    y += Inches(0.35 + len(items) * 0.22 + Inches(0.1).inches)

# 预期成果
add_rect(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.7), SECTION_BG)
add_text(slide, Inches(0.8), Inches(5.5), Inches(10), Inches(0.35),
         "预期成果", font_size=18, color=TITLE_BLUE, bold=True)

outcomes = [
    "一套完整的基于多智能体技术的Web无障碍智能检测系统 (W4Agent)",
    "与主流工具(axe-core、WAVE、Lighthouse)的对比实验数据与分析",
    "一篇本科毕业论文 (预计约3-4万字)",
    "系统部署文档及使用手册",
]
add_bullet_list(slide, Inches(0.7), Inches(5.9), Inches(11.5), Inches(1.2),
                outcomes, font_size=13, spacing=1.4, bullet_color=GREEN)

add_slide_number(slide, 13, TOTAL_SLIDES)


# ============================================================
# SLIDE 14: 致谢
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)

add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, ACCENT_BLUE)
add_rect(slide, Inches(0), Inches(3.3), SLIDE_W, Inches(0.04), ACCENT_TEAL)

add_text(slide, Inches(0), Inches(2.0), SLIDE_W, Inches(0.8),
         "感谢各位老师的指导与评审", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(3.7), SLIDE_W, Inches(0.5),
         "W4Agent — 基于智能体技术的Web无障碍智能检测系统", font_size=18, color=LIGHT_BLUE, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(4.5), SLIDE_W, Inches(0.5),
         "恳请各位老师批评指正", font_size=20, color=RGBColor(0xAA, 0xCC, 0xEE), alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(5.5), SLIDE_W, Inches(0.4),
         "答辩人：XXX    |    指导教师：XXX 教授", font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ── 保存 ──
output_path = "/Users/guangzhi/Desktop/W4Agent/中期答辩PPT.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
print(f"共 {TOTAL_SLIDES} 页")
